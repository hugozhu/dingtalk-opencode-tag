"""file — 文档/文件消息处理能力（custom 插件）(#40, #68)

群里发文件时，event-consume 下 content 形如
`[文件] <文件名> fileId: <fileId> 注意：如需下载使用dws drive download命令下载`。

**受控处理**（对齐 image 能力）：harness **主动**把文件下载到**临时目录**（不是项目
工作目录）、按类型解析、注入 agent、回复、用完删。避免"agent 自主用 bash 工具下载
文件到项目目录/执行 shell"的不可控 + 安全问题（见 #40）。

**精细化处理**（#68）：按文件类型分派到不同解析器，每类在独立 session 解析后包裹
提示词注入复用的主会话（保留多轮上下文）。

流程：
  1. on_inbound(kind=file)：防回环 → 去重 → 提交 handle_file。
  2. handle_file：提取 fileId + 文件名 → 按类型分派解析器 → 独立 session 解析 →
     包裹提示词注入复用主会话 → 回复发回群 → 删 tmpdir + 删临时 session。
  3. 下载失败 / 解析失败 / 不支持类型 → 明确告知用户，不静默、不硬塞乱码。

支持类型：
  - 文本：直接读前 N 字节（txt/md/csv/json/日志/代码等）
  - 图片：复用 image._recognize_via_serve（vision 识别，独立 session）
  - PDF：文本层提取优先；扫描/无文本层则逐页转图走 vision
  - Office：docx/xlsx/pptx 转纯文本/结构化文本
  - 视频：抽关键帧走 vision（可选音轨转写）
  - 其它：明确告知读不了

开关：CAP_FILE_ENABLED（默认开）。优先级 40（与 image 同级，先于 forward/text）。
"""

import os
import re
import shutil
import tempfile
import base64

from core.agent_common import _run_cli, log, submit_handler, find_serve_credentials, serve_request
from core.capabilities import Capability, register
from core.inbound import KIND_FILE
from core.brain import generate_reply, register_session as _register_textreply_sid
from core.replier import send_reply

# 从 content 提取 fileId 和文件名
# 格式：[文件] <文件名> fileId: <fileId> 注意：...
_RE_FILE_ID = re.compile(r"fileId:\s*(\S+)")
_RE_FILE_NAME = re.compile(r"\[文件\]\s*(.+?)\s+fileId:")

# 读取文本正文的字节上限（防超大文件撑爆 prompt）
_FILE_MAX_BYTES = int(os.environ.get("CAP_FILE_MAX_BYTES", "16384"))

# 文本类文件后缀
_TEXT_EXTS = {
    ".txt", ".md", ".markdown", ".csv", ".tsv", ".json", ".jsonl", ".log",
    ".yaml", ".yml", ".xml", ".ini", ".conf", ".cfg", ".toml", ".py", ".js",
    ".ts", ".sh", ".go", ".java", ".c", ".cpp", ".h", ".rs", ".rb", ".php",
    ".sql", ".html", ".css", ".env", ".properties", ".gitignore",
}

# 图片类后缀（走 vision 识别）
_IMAGE_EXTS = {".jpg", ".jpeg", ".png", ".gif", ".bmp", ".webp", ".svg"}

# PDF 后缀
_PDF_EXTS = {".pdf"}

# Office 文档后缀
_OFFICE_EXTS = {".docx", ".xlsx", ".pptx", ".doc", ".xls", ".ppt"}

# 视频后缀
_VIDEO_EXTS = {".mp4", ".avi", ".mov", ".mkv", ".flv", ".wmv", ".webm"}

# 防回环 + msgId 去重由 core 声明式处理（见 Capability(loop_guard/dedup)）。

# Vision 模型配置（复用 image 能力的配置）
_VISION_MODEL = os.environ.get("AGENT_VISION_MODEL", "")
_VISION_TIMEOUT = int(os.environ.get("CAP_IMAGE_VISION_TIMEOUT", "90"))

# 各类型 prompt footer（可由环境变量覆盖）
_FILE_TEXT_PROMPT_FOOTER = os.environ.get(
    "CAP_FILE_TEXT_PROMPT_FOOTER",
    "以上是用户发送的文件内容（由系统下载并读取前 {max_bytes} 字节，可能已截断）。\n"
    "请理解文件的完整语境（主题、关键信息、数据/代码细节），结合用户随文件的说明（若有），\n"
    "对用户的意图做出有帮助的回应（该答疑答疑、该归纳归纳、该指出问题指出问题）。",
).format(max_bytes=_FILE_MAX_BYTES)

_FILE_IMAGE_PROMPT_FOOTER = os.environ.get(
    "CAP_FILE_IMAGE_PROMPT_FOOTER",
    "以上「图片识别内容」由多模态模型从用户发送的图片文件中提取/描述得到。\n"
    "你本身看不到原图，但可根据识别内容理解用户意图。\n"
    "请结合用户随图的说明（若有），对用户发送这张图片的意图做出有帮助的回应。",
)

_FILE_PDF_PROMPT_FOOTER = os.environ.get(
    "CAP_FILE_PDF_PROMPT_FOOTER",
    "以上是从用户发送的 PDF 文件中提取的内容（由系统解析得到，你看不到原件）。\n"
    "请理解文档的完整语境，结合用户说明（若有），做出有帮助的回应。",
)

_FILE_OFFICE_PROMPT_FOOTER = os.environ.get(
    "CAP_FILE_OFFICE_PROMPT_FOOTER",
    "以上是从用户发送的 Office 文档中提取的内容（由系统解析得到，你看不到原件）。\n"
    "请理解文档的完整语境，结合用户说明（若有），做出有帮助的回应。",
)

_FILE_VIDEO_PROMPT_FOOTER = os.environ.get(
    "CAP_FILE_VIDEO_PROMPT_FOOTER",
    "以上是从用户发送的视频文件中提取的关键帧内容（由系统抽帧并识别得到，你看不到原视频）。\n"
    "请理解视频的内容和意图，结合用户说明（若有），做出有帮助的回应。",
)

# 各类型解析开关（默认全开，缺依赖时优雅降级）
_IMAGE_PARSER_ENABLED = os.environ.get("CAP_FILE_IMAGE_ENABLED", "1").lower() in ("1", "true", "yes", "on")
_PDF_PARSER_ENABLED = os.environ.get("CAP_FILE_PDF_ENABLED", "1").lower() in ("1", "true", "yes", "on")
_OFFICE_PARSER_ENABLED = os.environ.get("CAP_FILE_OFFICE_ENABLED", "1").lower() in ("1", "true", "yes", "on")
_VIDEO_PARSER_ENABLED = os.environ.get("CAP_FILE_VIDEO_ENABLED", "1").lower() in ("1", "true", "yes", "on")

# Vision 识别 prompt（与 image 能力一致）
_VISION_PROMPT = os.environ.get(
    "CAP_FILE_VISION_PROMPT",
    "请逐字提取这张图片中的所有文字内容（保持原始顺序、换行、标点，不要省略或总结）。"
    "如果图片中没有文字，则客观描述图片内容（场景、物体、UI 元素、图表数据、颜色等），"
    "不要做主观总结或解读。",
)


def _download_file(file_id):
    """drive download 到临时目录（**非项目目录**），返回本地路径或 None。"""
    tmp_dir = tempfile.mkdtemp(prefix="agent_file_")
    rc, _ = _run_cli([
        "drive", "download",
        "--node", file_id,
        "--output", tmp_dir + "/",
    ], timeout=60)
    if rc != 0:
        log(f"file: 下载失败 rc={rc} fileId={file_id[:16]}")
        shutil.rmtree(tmp_dir, ignore_errors=True)
        return None
    for name in os.listdir(tmp_dir):
        return os.path.join(tmp_dir, name)
    log(f"file: 下载目录为空 fileId={file_id[:16]}")
    shutil.rmtree(tmp_dir, ignore_errors=True)
    return None


def _classify_file(filename):
    """按文件名后缀分类，返回类型标识：text/image/pdf/office/video/unknown。"""
    _, ext = os.path.splitext((filename or "").lower())
    if ext in _TEXT_EXTS:
        return "text"
    if ext in _IMAGE_EXTS:
        return "image"
    if ext in _PDF_EXTS:
        return "pdf"
    if ext in _OFFICE_EXTS:
        return "office"
    if ext in _VIDEO_EXTS:
        return "video"
    return "unknown"


def _split_model(model):
    """'provider/model' → (providerID, modelID)。无 '/' 时 provider 空。"""
    if "/" in (model or ""):
        p, _, m = model.partition("/")
        return p, m
    return "", (model or "")


def _recognize_via_serve(img_bytes, mime="image/png"):
    """经 opencode serve 用 AGENT_VISION_MODEL 识别图片，返回描述文本或 ""。

    建临时 session → POST 一条含 file(image) + text 的 user 消息（阻塞到该轮完成）→
    取 assistant 文本 → best-effort 删 session。免外部 proxy。
    """
    if not _VISION_MODEL:
        return ""
    pid, port, pwd = find_serve_credentials()
    if not port:
        log("file: serve 凭据缺失，serve 识别不可用")
        return ""
    provider, model_id = _split_model(_VISION_MODEL)
    b64 = base64.b64encode(img_bytes).decode()
    data_url = f"data:{mime};base64,{b64}"

    sid = None
    try:
        created = serve_request("POST", "/session", {"title": "agent-file-vision"},
                                timeout=10, port=port, pwd=pwd)
        sid = (created or {}).get("id")
        if not sid:
            return ""
        # 登记为 brain 抑制名单，避免 vision session 的 SSE 事件触发业务通知刷屏
        _register_textreply_sid(sid)
        d = serve_request("POST", f"/session/{sid}/message", {
            "model": {"providerID": provider, "modelID": model_id},
            "parts": [
                {"type": "file", "mime": mime, "filename": "image", "url": data_url},
                {"type": "text", "text": _VISION_PROMPT},
            ],
        }, timeout=_VISION_TIMEOUT, port=port, pwd=pwd) or {}
        desc = "".join(
            p.get("text", "") for p in d.get("parts", []) if p.get("type") == "text"
        ).strip()
        if desc:
            log(f"file: serve 识别成功 model={_VISION_MODEL} desc_len={len(desc)}")
        return desc
    except Exception as e:
        log(f"file: serve 识别失败 model={_VISION_MODEL} err={e}")
        return ""
    finally:
        if sid:
            try:
                serve_request("DELETE", f"/session/{sid}", timeout=6, port=port, pwd=pwd)
            except Exception:
                pass


# ============================================================================
# 各类型解析器（返回 (content_text, success)）
# ============================================================================

def _parse_text(path):
    """读文件前 N 字节文本，返回 (text, success)。"""
    try:
        with open(path, "r", encoding="utf-8", errors="replace") as f:
            data = f.read(_FILE_MAX_BYTES + 1)
        truncated = len(data) > _FILE_MAX_BYTES
        text = data[:_FILE_MAX_BYTES]
        if truncated:
            text += f"\n\n（文件过长，仅读取前 {_FILE_MAX_BYTES} 字节）"
        return text, True
    except Exception as e:
        log(f"file: 文本读取失败 {e}")
        return "", False


def _parse_image(path):
    """图片走 vision 识别，返回 (desc, success)。"""
    if not _IMAGE_PARSER_ENABLED:
        return "", False
    mime = "image/jpeg" if path.lower().endswith((".jpg", ".jpeg")) else "image/png"
    try:
        with open(path, "rb") as f:
            img_bytes = f.read()
        desc = _recognize_via_serve(img_bytes, mime=mime)
        return desc, bool(desc)
    except Exception as e:
        log(f"file: 图片识别失败 {e}")
        return "", False


def _parse_pdf(path):
    """PDF 文本层提取；扫描版则逐页转图走 vision，返回 (text, success)。"""
    if not _PDF_PARSER_ENABLED:
        return "", False
    try:
        # 尝试导入 pdfplumber（文本层提取）
        import pdfplumber
        with pdfplumber.open(path) as pdf:
            texts = []
            for i, page in enumerate(pdf.pages[:20]):  # 最多前 20 页防超时
                text = page.extract_text() or ""
                if text.strip():
                    texts.append(f"--- 第 {i+1} 页 ---\n{text}")
            if texts:
                return "\n\n".join(texts), True
            # 无文本层，走 OCR 回退
            log(f"file: PDF 无文本层，尝试 OCR")
            return _parse_pdf_ocr(path)
    except ImportError:
        log("file: PDF 解析器不可用（缺 pdfplumber），尝试 OCR 回退")
        return _parse_pdf_ocr(path)
    except Exception as e:
        log(f"file: PDF 解析失败 {e}")
        return "", False


def _parse_pdf_ocr(path):
    """PDF 逐页转图走 vision OCR，返回 (text, success)。"""
    try:
        from pdf2image import convert_from_path
        images = convert_from_path(path, first_page=1, last_page=10)  # 最多前 10 页
        texts = []
        for i, img in enumerate(images):
            # 转 bytes
            import io
            buf = io.BytesIO()
            img.save(buf, format="PNG")
            img_bytes = buf.getvalue()
            desc = _recognize_via_serve(img_bytes, mime="image/png")
            if desc:
                texts.append(f"--- 第 {i+1} 页 ---\n{desc}")
        if texts:
            return "\n\n".join(texts), True
        return "", False
    except ImportError:
        log("file: PDF OCR 不可用（缺 pdf2image）")
        return "", False
    except Exception as e:
        log(f"file: PDF OCR 失败 {e}")
        return "", False


def _parse_office(path):
    """Office 文档转纯文本/结构化文本，返回 (text, success)。"""
    if not _OFFICE_PARSER_ENABLED:
        return "", False
    ext = os.path.splitext(path.lower())[1]
    try:
        if ext == ".docx":
            return _parse_docx(path)
        elif ext == ".xlsx":
            return _parse_xlsx(path)
        elif ext == ".pptx":
            return _parse_pptx(path)
        else:
            # 老格式 doc/xls/ppt 需要 libreoffice 转换，暂不支持
            log(f"file: 老格式 Office 文档暂不支持 {ext}")
            return "", False
    except Exception as e:
        log(f"file: Office 解析失败 {e}")
        return "", False


def _parse_docx(path):
    """docx 提取文本，返回 (text, success)。"""
    try:
        from docx import Document
        doc = Document(path)
        texts = [para.text for para in doc.paragraphs if para.text.strip()]
        return "\n\n".join(texts), True
    except ImportError:
        log("file: docx 解析器不可用（缺 python-docx）")
        return "", False
    except Exception as e:
        log(f"file: docx 解析失败 {e}")
        return "", False


def _parse_xlsx(path):
    """xlsx 提取表格（保留行列结构），返回 (text, success)。"""
    try:
        from openpyxl import load_workbook
        wb = load_workbook(path, read_only=True, data_only=True)
        texts = []
        for sheet_name in wb.sheetnames[:5]:  # 最多前 5 个 sheet
            sheet = wb[sheet_name]
            texts.append(f"--- Sheet: {sheet_name} ---")
            rows = []
            for row in list(sheet.iter_rows(values_only=True))[:100]:  # 最多前 100 行
                row_text = "\t".join(str(cell) if cell is not None else "" for cell in row)
                if row_text.strip():
                    rows.append(row_text)
            texts.append("\n".join(rows))
        return "\n\n".join(texts), True
    except ImportError:
        log("file: xlsx 解析器不可用（缺 openpyxl）")
        return "", False
    except Exception as e:
        log(f"file: xlsx 解析失败 {e}")
        return "", False


def _parse_pptx(path):
    """pptx 提取文本（每页标题 + 正文），返回 (text, success)。"""
    try:
        from pptx import Presentation
        prs = Presentation(path)
        texts = []
        for i, slide in enumerate(prs.slides[:20]):  # 最多前 20 页
            slide_texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text)
            if slide_texts:
                texts.append(f"--- 第 {i+1} 页 ---\n" + "\n".join(slide_texts))
        return "\n\n".join(texts), True
    except ImportError:
        log("file: pptx 解析器不可用（缺 python-pptx）")
        return "", False
    except Exception as e:
        log(f"file: pptx 解析失败 {e}")
        return "", False


def _parse_video(path):
    """视频抽关键帧走 vision，返回 (text, success)。"""
    if not _VIDEO_PARSER_ENABLED:
        return "", False
    try:
        import cv2
        cap = cv2.VideoCapture(path)
        if not cap.isOpened():
            return "", False
        fps = cap.get(cv2.CAP_PROP_FPS) or 25
        frame_count = int(cap.get(cv2.CAP_PROP_FRAME_COUNT))
        # 抽 5 帧（开头、1/4、1/2、3/4、结尾）
        positions = [0, frame_count // 4, frame_count // 2, frame_count * 3 // 4, frame_count - 1]
        texts = []
        for i, pos in enumerate(positions):
            cap.set(cv2.CAP_PROP_POS_FRAMES, pos)
            ret, frame = cap.read()
            if not ret:
                continue
            # 转 PNG bytes
            import io
            success, buf = cv2.imencode(".png", frame)
            if not success:
                continue
            img_bytes = buf.tobytes()
            desc = _recognize_via_serve(img_bytes, mime="image/png")
            if desc:
                timestamp = pos / fps if fps > 0 else 0
                texts.append(f"--- 时间 {timestamp:.1f}s ---\n{desc}")
        cap.release()
        if texts:
            return "\n\n".join(texts), True
        return "", False
    except ImportError:
        log("file: 视频解析器不可用（缺 opencv-python）")
        return "", False
    except Exception as e:
        log(f"file: 视频解析失败 {e}")
        return "", False


# ============================================================================
# 主处理逻辑
# ============================================================================

def handle_file(user, text, msg_id, conv_id, conv_type):
    """提取 fileId+文件名 → 分类 → 解析 → 包裹提示词 → 注入复用主会话 → 回复 → 删 tmpdir。"""
    fid_m = _RE_FILE_ID.search(text or "")
    if not fid_m:
        log(f"file: 未提取到 fileId msgId={msg_id[:24]}")
        return
    file_id = fid_m.group(1)
    name_m = _RE_FILE_NAME.search(text or "")
    filename = name_m.group(1).strip() if name_m else "（未知文件名）"

    file_type = _classify_file(filename)
    log(f"file: msgId={msg_id[:24]} filename={filename!r} type={file_type}")

    path = _download_file(file_id)
    if not path:
        send_reply(conv_id, conv_type, f"抱歉，文件「{filename}」我没能下载下来，能再发一次吗？")
        return

    tmp_dir = os.path.dirname(path)
    content = None
    success = False
    footer = ""

    try:
        if file_type == "text":
            content, success = _parse_text(path)
            footer = _FILE_TEXT_PROMPT_FOOTER
        elif file_type == "image":
            content, success = _parse_image(path)
            footer = _FILE_IMAGE_PROMPT_FOOTER
        elif file_type == "pdf":
            content, success = _parse_pdf(path)
            footer = _FILE_PDF_PROMPT_FOOTER
        elif file_type == "office":
            content, success = _parse_office(path)
            footer = _FILE_OFFICE_PROMPT_FOOTER
        elif file_type == "video":
            content, success = _parse_video(path)
            footer = _FILE_VIDEO_PROMPT_FOOTER
        else:
            # 未知类型
            send_reply(conv_id, conv_type,
                       f"收到文件「{filename}」，但它看起来不是我能处理的文件类型。\n"
                       f"我可以处理：文本文件（txt/md/csv/json/日志/代码等）、图片、PDF、Office 文档（docx/xlsx/pptx）、视频。")
            return
    finally:
        shutil.rmtree(tmp_dir, ignore_errors=True)  # 用完即删，不留临时文件

    if not success or not content:
        send_reply(conv_id, conv_type, f"抱歉，文件「{filename}」的内容我解析失败了。")
        return

    log(f"file: msgId={msg_id[:24]} 解析成功 type={file_type} content_len={len(content)}")

    # 参考 image 能力：结构化呈现文件信息（用户+文件名+内容+任务指令）
    parts = [
        f"用户 {user} 发送了一个文件：{filename}",
        "",
        "【文件内容】",
        "```",
        content,
        "```",
        "",
        footer,
    ]
    prompt = "\n".join(parts)

    # 注入复用的主会话（对齐 image.py:196，raw=True + ctx 传递会话信息）
    reply = generate_reply(user, prompt, raw=True, ctx={
        "conv_id": conv_id, "conv_type": conv_type, "msg_id": msg_id, "user": user,
    })
    if reply:
        send_reply(conv_id, conv_type, reply)
    else:
        log(f"file: 大脑无回复 msgId={msg_id[:24]}")


def on_inbound(msg):
    """文件消息入站：提交 handle_file。返回 True=已消费。

    防回环 + msgId 去重由 core dispatch_inbound（loop_guard/dedup 声明）处理。
    """
    submit_handler(handle_file, msg.user, msg.text, msg.msg_id, msg.conv_id, msg.conv_type)
    return True


CAPABILITY = Capability(
    name="file",
    on_inbound=on_inbound,
    handles_kinds={KIND_FILE},
    priority=40,             # 与 image 同级，先于 forward(50)/text(100)
    default_enabled=True,
    loop_guard=True,         # core 统一防回环
    dedup=True,              # core 统一 msgId 去重
)
register(CAPABILITY)
